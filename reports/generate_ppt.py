from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
import os
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import tempfile
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import io
import tempfile

def clean_text(text):
    """Remove all markdown formatting and clean text aggressively"""
    if not text:
        return ""
    
    # Convert to string
    text = str(text)
    
    # Remove all asterisks (multiple passes to catch nested ones)
    while "**" in text:
        text = text.replace("**", "")
    while "*" in text:
        text = text.replace("*", "")
    
    # Remove other markdown characters
    text = text.replace("#", "")
    text = text.replace("_", "")
    text = text.replace("`", "")
    text = text.replace("~", "")
    text = text.replace("|", "")
    
    # Remove markdown links [text](url) -> text
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    
    # Remove HTML tags if any
    text = re.sub(r'<[^>]+>', '', text)
    
    # Clean up extra spaces and newlines
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n+', ' ', text)
    
    # Remove leading/trailing punctuation that might be left over
    text = text.strip(' .,;:!?')
    
    return text.strip()

def generate_ppt(company, insights):
    prs = Presentation()

    # Set slide width and height (standard 16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Color scheme
    PRIMARY_BLUE = RGBColor(30, 58, 138)
    SECONDARY_BLUE = RGBColor(59, 130, 246)
    DARK_GRAY = RGBColor(31, 41, 55)
    LIGHT_GRAY = RGBColor(100, 116, 139)
    BACKGROUND_GRAY = RGBColor(249, 250, 251)

    def add_title_slide():
        """Create professional title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background rectangle
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = PRIMARY_BLUE
        rect.line.fill.background()
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = "Intelligent Business Strategy Orchestrator"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(42)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Executive Decision Overview"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = LIGHT_GRAY
        subtitle_para.alignment = PP_ALIGN.CENTER

    def add_metric_slide(title, metrics):
        """Create a slide with key metrics in boxes"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE
        title_para.alignment = PP_ALIGN.LEFT
        
        # Add metrics in boxes (left side - smaller to make room for chart)
        box_width = Inches(4.0)
        box_height = Inches(1.0)
        spacing = Inches(0.25)
        start_y = Inches(1.6)
        
        for i, (label, value) in enumerate(metrics.items()):
            row = i // 2
            col = i % 2
            left_pos = Inches(0.5) + col * (box_width + spacing)
            top_pos = start_y + row * (box_height + spacing)
            
            # Background box
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, box_width, box_height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = BACKGROUND_GRAY
            rect.line.color.rgb = SECONDARY_BLUE
            rect.line.width = Pt(2)
            
            # Label
            label_box = slide.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), 
                                                 box_width - Inches(0.4), Inches(0.4))
            label_frame = label_box.text_frame
            label_frame.text = str(label)
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = LIGHT_GRAY
            label_para.alignment = PP_ALIGN.LEFT
            
            # Value
            value_box = slide.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.5), 
                                                 box_width - Inches(0.4), Inches(0.6))
            value_frame = value_box.text_frame
            value_frame.text = str(value)
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(28)
            value_para.font.bold = True
            value_para.font.color.rgb = PRIMARY_BLUE
            value_para.alignment = PP_ALIGN.LEFT
        
        # Add pie chart for Revenue vs Expenses (right side)
        try:
            revenue = company.get('revenue', 0)
            expenses = company.get('expenses', 0)
            if revenue > 0 or expenses > 0:
                # Create pie chart
                fig, ax = plt.subplots(figsize=(4, 4))
                colors_pie = ['#2563eb', '#dc2626']
                labels_pie = ['Revenue', 'Expenses']
                sizes = [revenue, expenses]
                
                # Only show if both values are positive
                if revenue > 0 and expenses > 0:
                    wedges, texts, autotexts = ax.pie(sizes, labels=labels_pie, colors=colors_pie, 
                                                     autopct='%1.1f%%', startangle=90,
                                                     textprops={'fontsize': 12, 'fontweight': 'bold'})
                    ax.set_title('Revenue vs Expenses', fontsize=14, fontweight='bold', pad=15)
                    
                    # Save to temporary file
                    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                    plt.savefig(temp_file.name, dpi=150, bbox_inches='tight', facecolor='white')
                    plt.close()
                    
                    # Insert chart into slide
                    chart_left = Inches(5.5)
                    chart_top = Inches(1.8)
                    chart_width = Inches(4.0)
                    chart_height = Inches(3.5)
                    slide.shapes.add_picture(temp_file.name, chart_left, chart_top, chart_width, chart_height)
                    
                    # Clean up
                    try:
                        os.unlink(temp_file.name)
                    except:
                        pass
        except Exception as e:
            pass  # Skip chart if there's an error

    def add_bullet_slide(title, items, max_items=5):
        """Create a clean bullet point slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE
        title_para.alignment = PP_ALIGN.LEFT
        
        # Underline
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(3), Inches(0.08))
        underline.fill.solid()
        underline.fill.fore_color.rgb = SECONDARY_BLUE
        underline.line.fill.background()
        
        # Content area (left side - smaller to make room for chart)
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.vertical_anchor = MSO_ANCHOR.TOP
        content_frame.clear()
        
        # Process items - clean aggressively
        if not isinstance(items, list):
            items = [str(items)]
        
        display_items = []
        for item in items[:max_items]:
            if item:
                # Clean the text
                cleaned = clean_text(str(item))
                if cleaned:
                    # Remove bullet symbols if already present
                    cleaned = re.sub(r'^[•\-\*\+\d+\.\)]\s*', '', cleaned)
                    cleaned = cleaned.strip()
                    if cleaned:
                        display_items.append(cleaned)
        
        # Add paragraphs
        for i, item_text in enumerate(display_items):
            if item_text:
                p = content_frame.add_paragraph()
                p.text = item_text
                p.level = 0
                p.font.size = Pt(22)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(18)
                p.alignment = PP_ALIGN.LEFT
                p.line_spacing = 1.3
                p.left_indent = Inches(0.3)
                p.first_line_indent = Inches(-0.3)
                
                if i == 0:
                    p.space_before = Pt(0)
                else:
                    p.space_before = Pt(10)
        
        # Add horizontal bar chart (right side)
        try:
            if display_items:
                fig, ax = plt.subplots(figsize=(4, 3.5))
                # Create priority values (higher priority = higher value)
                priorities = list(range(len(display_items), 0, -1))
                items_short = [item[:25] + "..." if len(item) > 25 else item for item in display_items]
                
                bars = ax.barh(items_short, priorities, color='#3b82f6', alpha=0.8, edgecolor='#1e40af', linewidth=1.5)
                ax.set_xlabel('Priority', fontsize=11, fontweight='bold')
                ax.set_title('Strategic Priorities', fontsize=13, fontweight='bold', pad=10)
                ax.set_xlim(0, len(display_items) + 1)
                ax.grid(True, alpha=0.3, axis='x', linestyle='--')
                ax.spines['top'].set_visible(False)
                ax.spines['right'].set_visible(False)
                
                # Add value labels on bars
                for i, (bar, priority) in enumerate(zip(bars, priorities)):
                    width = bar.get_width()
                    ax.text(width + 0.1, bar.get_y() + bar.get_height()/2, 
                           f'#{priority}', ha='left', va='center', fontweight='bold', fontsize=10)
                
                plt.tight_layout()
                
                # Save to temporary file
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                plt.savefig(temp_file.name, dpi=150, bbox_inches='tight', facecolor='white')
                plt.close()
                
                # Insert chart into slide
                chart_left = Inches(5.5)
                chart_top = Inches(1.8)
                chart_width = Inches(4.0)
                chart_height = Inches(4.5)
                slide.shapes.add_picture(temp_file.name, chart_left, chart_top, chart_width, chart_height)
                
                # Clean up
                try:
                    os.unlink(temp_file.name)
                except:
                    pass
        except Exception as e:
            pass  # Skip chart if there's an error

    def add_text_slide(title, text_content):
        """Create a slide with formatted text paragraphs"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE
        title_para.alignment = PP_ALIGN.LEFT
        
        # Underline
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(4.5), Inches(0.08))
        underline.fill.solid()
        underline.fill.fore_color.rgb = SECONDARY_BLUE
        underline.line.fill.background()
        
        # Content area
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.vertical_anchor = MSO_ANCHOR.TOP
        content_frame.margin_left = Inches(0)
        content_frame.margin_right = Inches(0)
        content_frame.margin_top = Inches(0)
        content_frame.margin_bottom = Inches(0)
        content_frame.clear()
        
        # Process text content - clean aggressively
        if isinstance(text_content, list):
            sentences = []
            for item in text_content:
                if item:
                    # Clean each item
                    cleaned = clean_text(str(item))
                    if cleaned:
                        # Split into sentences if it's long
                        if len(cleaned) > 150:
                            split_sentences = [s.strip() for s in re.split(r'[.!?]+', cleaned) if s.strip()]
                            sentences.extend(split_sentences)
                        else:
                            sentences.append(cleaned)
        else:
            # Split by sentences
            text = clean_text(str(text_content))
            sentences = [s.strip() for s in re.split(r'[.!?]+', text) if s.strip()]
        
        # Clean sentences again to ensure no asterisks
        cleaned_sentences = []
        for s in sentences:
            cleaned = clean_text(s)
            if cleaned and len(cleaned) > 5:  # Only add meaningful sentences
                cleaned_sentences.append(cleaned)
        
        # Add paragraphs
        for i, sentence in enumerate(cleaned_sentences[:7]):  # Limit to 7 sentences
            if sentence:
                # Ensure sentence ends properly
                if not sentence[-1] in ".!?":
                    sentence += "."
                
                # Final clean before adding
                sentence = clean_text(sentence)
                
                p = content_frame.add_paragraph()
                p.text = sentence
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(16)
                p.alignment = PP_ALIGN.LEFT
                p.line_spacing = 1.4
                p.left_indent = Inches(0)
                p.first_line_indent = Inches(0)
                
                if i == 0:
                    p.space_before = Pt(0)
                else:
                    p.space_before = Pt(10)

    # Title Slide
    add_title_slide()

    # Business Snapshot
    profit_loss = company['revenue'] - company['expenses']
    metrics = {
        "Revenue": f"₹{company['revenue']:,}",
        "Expenses": f"₹{company['expenses']:,}",
        "Profit/Loss": f"₹{profit_loss:,}",
        "Risk Level": str(insights.get('risk', 'N/A'))
    }
    add_metric_slide("Business Snapshot", metrics)

    # Strategic Direction - clean the data first
    strategy_items = insights.get("strategy_focus", [])
    if not isinstance(strategy_items, list):
        strategy_items = [str(strategy_items)]
    
    # Clean all strategy items
    cleaned_strategy = []
    for item in strategy_items:
        if item:
            cleaned = clean_text(str(item))
            if cleaned:
                cleaned_strategy.append(cleaned)
    
    add_bullet_slide("Strategic Direction", cleaned_strategy, max_items=5)

    # Executive Direction & Priorities - clean the data first
    ceo_items = insights.get("ceo_summary", [])
    ceo_list = []
    
    if isinstance(ceo_items, list):
        for item in ceo_items:
            if item:
                # Clean immediately
                cleaned_item = clean_text(str(item))
                if cleaned_item:
                    # Split by sentences if it's a long string
                    if len(cleaned_item) > 100:
                        sentences = [s.strip() for s in re.split(r'[.!?]+', cleaned_item) if s.strip()]
                        for sent in sentences:
                            cleaned_sent = clean_text(sent)
                            if cleaned_sent:
                                ceo_list.append(cleaned_sent)
                    else:
                        ceo_list.append(cleaned_item)
    elif ceo_items:
        cleaned = clean_text(str(ceo_items))
        if cleaned:
            ceo_list.append(cleaned)
    
    add_text_slide("Executive Direction & Priorities", ceo_list)

    path = "reports/Executive_Presentation.pptx"
    
    # Delete old file if it exists to avoid caching issues
    if os.path.exists(path):
        try:
            os.remove(path)
        except:
            pass
    
    prs.save(path)
    return path
